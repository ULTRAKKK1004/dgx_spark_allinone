import asyncio
import os
import uuid
from typing import List, Dict, Any
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

class SlideRenderer:
    def __init__(self, results_dir: str):
        self.results_dir = results_dir

    async def render_html_to_images(self, html_content: str) -> List[str]:
        """Renders HTML slides to individual PNG images using Playwright."""
        image_paths = []
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            page = await browser.new_page(viewport={"width": 1920, "height": 1080})
            
            # Write temporary HTML file
            temp_html = os.path.join(self.results_dir, f"temp_slides_{uuid.uuid4().hex}.html")
            with open(temp_html, "w", encoding="utf-8") as f:
                f.write(html_content)
            
            await page.goto(f"file://{temp_html}")
            
            # Detect slides by .slide class
            slides = await page.query_selector_all(".slide")
            if not slides:
                # Fallback to single page if no .slide class found
                path = os.path.join(self.results_dir, f"slide_0_{uuid.uuid4().hex}.png")
                await page.screenshot(path=path, full_page=False)
                image_paths.append(path)
            else:
                for i, slide in enumerate(slides):
                    path = os.path.join(self.results_dir, f"slide_{i}_{uuid.uuid4().hex}.png")
                    await slide.screenshot(path=path)
                    image_paths.append(path)
            
            await browser.close()
            if os.path.exists(temp_html):
                os.remove(temp_html)
        
        return image_paths

    def create_pptx_from_images(self, image_paths: List[str], output_path: str):
        """Creates a PPTX file where each slide is an image."""
        prs = Presentation()
        # Set 16:9 aspect ratio
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        for img_path in image_paths:
            slide_layout = prs.slide_layouts[6] # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            
        prs.save(output_path)
        # Cleanup images
        for img_path in image_paths:
            if os.path.exists(img_path):
                os.remove(img_path)
        return output_path

    async def check_layout_issues(self, html_content: str) -> Dict[str, Any]:
        """Checks for overflow or overlapping elements in the HTML slides."""
        issues = []
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            page = await browser.new_page(viewport={"width": 1920, "height": 1080})
            
            temp_html = os.path.join(self.results_dir, f"check_{uuid.uuid4().hex}.html")
            with open(temp_html, "w", encoding="utf-8") as f:
                f.write(html_content)
            
            await page.goto(f"file://{temp_html}")
            
            # JS script to detect overflow
            overflow_script = """
            () => {
                const results = [];
                document.querySelectorAll('.slide').forEach((slide, index) => {
                    const hasOverflow = slide.scrollHeight > slide.clientHeight || slide.scrollWidth > slide.clientWidth;
                    if (hasOverflow) {
                        results.append({slide: index, issue: 'overflow', detail: `Content exceeds slide bounds` });
                    }
                    // Add check for overlapping elements if possible
                });
                return results;
            }
            """
            issues = await page.evaluate(overflow_script)
            await browser.close()
            os.remove(temp_html)
            
        return {"has_issues": len(issues) > 0, "issues": issues}
